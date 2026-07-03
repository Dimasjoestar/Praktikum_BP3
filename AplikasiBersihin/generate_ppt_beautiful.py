from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
# Set 16:9 aspect ratio
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors matching the App Theme
teal = RGBColor(0x00, 0x96, 0x88)
dark_teal = RGBColor(0x00, 0x79, 0x6B)
bg_light = RGBColor(0xF5, 0xF9, 0xF9)
text_dark = RGBColor(0x1A, 0x2B, 0x2B)
white = RGBColor(0xFF, 0xFF, 0xFF)
accent = RGBColor(0x4D, 0xB6, 0xAC)
gold = RGBColor(0xFF, 0xD5, 0x4F)

def add_title_slide(prs, title_text, subtitle_text, author_text):
    slide_layout = prs.slide_layouts[6] # Blank slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Background full Teal
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = teal
    bg.line.fill.background()
    
    # Decorative shapes mimicking modern UI
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-2), Inches(8), Inches(8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.line.fill.background()
    
    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(5), Inches(4), Inches(4))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = dark_teal
    circle2.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.3), Inches(11.33), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = white
    
    p2 = tf.add_paragraph()
    p2.text = subtitle_text
    p2.font.size = Pt(28)
    p2.font.color.rgb = white
    p2.space_before = Pt(10)
    
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(11.33), Inches(1))
    tf2 = txBox2.text_frame
    p3 = tf2.paragraphs[0]
    p3.text = f"Oleh: {author_text}"
    p3.font.size = Pt(22)
    p3.font.color.rgb = gold
    p3.font.bold = True

def add_content_slide(prs, title_text, items):
    slide_layout = prs.slide_layouts[6] # Blank slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Light background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_light
    bg.line.fill.background()
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header.fill.solid()
    header.fill.fore_color.rgb = teal
    header.line.fill.background()
    
    # App Logo Text in top right
    logo_box = slide.shapes.add_textbox(Inches(10), Inches(0.3), Inches(3), Inches(0.8))
    logo_tf = logo_box.text_frame
    lp = logo_tf.paragraphs[0]
    lp.text = "Bersih.in"
    lp.font.size = Pt(24)
    lp.font.bold = True
    lp.font.color.rgb = white
    lp.alignment = PP_ALIGN.RIGHT
    
    # Title Text
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = white
    
    # Main Content Area (White Card)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.2))
    card.fill.solid()
    card.fill.fore_color.rgb = white
    card.line.fill.background()
    card.adjustments[0] = 0.03 # slight rounded corners
    
    # Bullets
    content_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11), Inches(4.5))
    ctf = content_box.text_frame
    ctf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            cp = ctf.paragraphs[0]
        else:
            cp = ctf.add_paragraph()
        # Bold the keyword before colon if present
        if ":" in item:
            parts = item.split(":", 1)
            run = cp.add_run()
            run.text = f"• {parts[0]}:"
            run.font.bold = True
            run.font.color.rgb = teal
            run2 = cp.add_run()
            run2.text = parts[1]
            run2.font.color.rgb = text_dark
        else:
            run = cp.add_run()
            run.text = f"• {item}"
            run.font.color.rgb = text_dark
        
        cp.font.size = Pt(26)
        cp.space_after = Pt(24)

def add_outro_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_light
    bg.line.fill.background()
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2), Inches(10.33), Inches(3.5))
    card.fill.solid()
    card.fill.fore_color.rgb = teal
    card.line.fill.background()
    card.adjustments[0] = 0.05
    
    txBox = slide.shapes.add_textbox(Inches(1.8), Inches(2.5), Inches(9.7), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "\"Sistem ini menciptakan ekosistem manajemen laundry yang terintegrasi, memberikan transparansi bagi pelanggan, kemudahan operasional bagi petugas, dan kontrol bisnis yang lebih baik bagi pengelola.\""
    p.font.size = Pt(28)
    p.font.italic = True
    p.font.color.rgb = white
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.2
    
    # Terima Kasih
    tkBox = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(11.33), Inches(1))
    tk_tf = tkBox.text_frame
    tk_p = tk_tf.paragraphs[0]
    tk_p.text = "TERIMA KASIH"
    tk_p.font.size = Pt(36)
    tk_p.font.bold = True
    tk_p.font.color.rgb = dark_teal
    tk_p.alignment = PP_ALIGN.CENTER

# ----------------- SLIDE GENERATION -----------------

add_title_slide(prs, "BERSIH.IN", "Sistem Layanan & Manajemen Laundry Premium\nIntegrasi Mobile Android & Panel Admin Terpadu", "Dimas Nurdiansyah")

add_content_slide(prs, "Latar Belakang & Masalah", [
    "Kurangnya digitalisasi pada sistem pemesanan dan pelacakan cucian.",
    "Pelanggan sulit memantau status pengerjaan pakaian secara real-time.",
    "Pengelola dan petugas masih merekap data pesanan secara manual yang rentan hilang atau selisih.",
    "Kurangnya engagement pelanggan karena tidak ada sistem reward yang menahan pelanggan (retention)."
])

add_content_slide(prs, "Solusi yang Ditawarkan", [
    "Aplikasi Android: Memudahkan eksplorasi layanan, pemesanan tiket pickup, dan sistem poin langsung dari genggaman.",
    "Backend Admin Panel: Dashboard untuk pemilik usaha mengelola data transaksi, reward, dan layanan dengan aman.",
    "Integrasi Logistik: Sistem pelacakan antar-jemput pakaian secara real-time yang terhubung dengan petugas."
])

add_content_slide(prs, "Teknologi yang Digunakan", [
    "Mobile Frontend: Kotlin Android, Android Jetpack Compose, Material Design 3.",
    "Arsitektur Data: Room Database (SQLite), DataStore, dan StateFlow / MVVM.",
    "Infrastruktur Bisnis: Sistem Multi-Role (Admin, Petugas, Pelanggan) dengan sistem manajemen state yang reaktif."
])

add_content_slide(prs, "Arsitektur Role & Hak Akses", [
    "Admin Utama: Kontrol sistem pusat. Kelola layanan, harga, promo, reward, dan pantau statistik pendapatan bulanan.",
    "Petugas (Officer): Otoritas lapangan. Verifikasi penjemputan, memproses status (Cuci/Setrika), dan selesaikan pesanan.",
    "Pelanggan (Customer): End-user aplikasi. Membuat order baru, menerima update real-time, dan klaim promo."
])

add_content_slide(prs, "Fitur: Sistem Order & Real-Time Tracking", [
    "Alur Transparan: Status berjenjang dari Menunggu Pickup -> Sedang Diproses -> Selesai.",
    "Otomatisasi Kalkulasi: Sistem otomatis menghitung total harga berdasarkan layanan dan tipe pengantaran.",
    "Riwayat Detil: Pengguna dapat mengakses riwayat order terdahulu sebagai bukti transaksi."
])

add_content_slide(prs, "Fitur: Sistem Reward & Gamifikasi", [
    "Sistem XP & Level: Pengguna otomatis mendapatkan XP di setiap transaksi sukses untuk naik level (Bronze, Silver, Gold).",
    "Penukaran Promo: XP dapat ditukar menjadi voucher diskon / promo eksklusif.",
    "Manajemen Dinamis: Admin dapat menambah atau mencabut promo yang tersedia kapan saja."
])

add_content_slide(prs, "Fitur Baru: Pembaruan & Optimalisasi", [
    "Ekspor Laporan PDF/CSV: Admin dapat mengunduh rekap keuangan bulanan langsung ke folder Downloads lokal dengan tombol ekspor khusus.",
    "Penyimpanan MediaStore: Sistem penyimpanan modern Android 11+ bebas izin runtime, membuat file langsung terlihat di File Manager.",
    "Pull-to-Refresh: Penambahan interaksi gestur tarik-bawah reaktif pada Dashboard Pelanggan, Petugas, dan Admin.",
    "Kueri Pendapatan Akurat: Sinkronisasi nominal pendapatan di dashboard Admin (hanya menghitung pesanan selesai/DONE).",
    "Ikon Aplikasi Kustom: Launcher icon toska gelembung air premium baru yang mendukung format adaptif (kotak/bulat)."
])

add_outro_slide(prs)

out_path = os.path.join(os.getcwd(), 'Presentasi_Aplikasi_Bersihin_Beautiful.pptx')
prs.save(out_path)
print(f"Presentation saved successfully to: {out_path}")
