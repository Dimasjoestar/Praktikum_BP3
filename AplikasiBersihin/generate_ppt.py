from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def add_title_slide(prs, title_text, subtitle_text, author_text):
    slide_layout = prs.slide_layouts[0] # Title slide
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = title_text
    subtitle.text = f"{subtitle_text}\n\nOleh: {author_text}"
    
    # Custom colors
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(0x00, 0x80, 0x80) # Teal
            
def add_bullet_slide(prs, title_text, bullets):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = title_text
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(0x00, 0x80, 0x80)

    tf = body.text_frame
    tf.text = bullets[0]
    for bullet in bullets[1:]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0

prs = Presentation()

# Slide 1: Judul
add_title_slide(prs, "BERSIH.IN", "Sistem Layanan & Manajemen Laundry Premium\nAplikasi Mobile Android untuk Pelanggan, Petugas, dan Admin", "Dimas Nurdiansyah")

# Slide 2: Latar Belakang & Masalah
add_bullet_slide(prs, "Latar Belakang & Masalah", [
    "Sulitnya memantau status pengerjaan pakaian secara real-time.",
    "Sistem pencatatan pesanan laundry di banyak tempat masih manual dan rentan hilang.",
    "Kurangnya engagement pelanggan karena tidak ada sistem promosi & reward yang jelas."
])

# Slide 3: Solusi yang Ditawarkan
add_bullet_slide(prs, "Solusi yang Ditawarkan", [
    "Aplikasi Pelanggan (Android): Memudahkan pelanggan membuat pesanan (pickup/antar sendiri), melacak status, dan mengumpulkan XP (Gamifikasi).",
    "Aplikasi Petugas: Membantu petugas mengelola penjemputan barang dan memperbarui status cucian.",
    "Portal Admin: Dashboard terpusat untuk memantau performa bisnis, mengatur layanan, promo, dan seluruh pesanan secara efisien."
])

# Slide 4: Teknologi yang Digunakan
add_bullet_slide(prs, "Teknologi yang Digunakan", [
    "Mobile UI: Kotlin, Android Jetpack Compose, Material Design 3.",
    "Local Database: Room Database (SQLite) & DataStore.",
    "Arsitektur: MVVM (Model-View-ViewModel) Pattern."
])

# Slide 5: Arsitektur Role & Hak Akses
add_bullet_slide(prs, "Arsitektur Role & Hak Akses", [
    "Admin: Kontrol penuh atas layanan, harga, promo, reward, dan laporan pendapatan.",
    "Petugas (Officer): Verifikasi pickup, penyesuaian pesanan, dan pembaruan status pengerjaan (Cuci/Setrika/Selesai).",
    "Pelanggan (Customer): Membuat order, melacak status, dan menukar poin/XP."
])

# Slide 6: Fitur: Pemantauan Status Real-Time
add_bullet_slide(prs, "Fitur: Pemantauan Status Real-Time", [
    "Pelanggan dapat melacak status pesanan secara real-time.",
    "Alur Transparan: Menunggu Pickup -> Diproses -> Selesai.",
    "Notifikasi status progres terstruktur pada antarmuka aplikasi."
])

# Slide 7: Fitur: Sistem Reward & Gamifikasi
add_bullet_slide(prs, "Fitur: Sistem Reward & Gamifikasi", [
    "Sistem Level: Pengguna mendapatkan XP setiap kali transaksi untuk naik level (Bronze, Silver, Gold).",
    "Penukaran Promo: XP dapat ditukar dengan voucher diskon khusus.",
    "Manajemen Dinamis: Reward dikelola secara dinamis dari sisi Admin."
])

# Slide 8: Kesimpulan
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Kesimpulan"
for paragraph in title.text_frame.paragraphs:
    for run in paragraph.runs:
        run.font.color.rgb = RGBColor(0x00, 0x80, 0x80)

body = slide.placeholders[1]
tf = body.text_frame
p = tf.paragraphs[0]
p.text = "\"Bersih.in menciptakan ekosistem manajemen laundry yang terintegrasi, memberikan transparansi bagi pelanggan, kemudahan operasional bagi petugas, dan kontrol bisnis yang lebih baik bagi pengelola.\""
# In python-pptx, font formatting on the first paragraph requires accessing its runs, but sometimes it doesn't have runs yet.
if not p.runs:
    run = p.add_run()
    run.text = p.text
    p.text = "" # Clear text
else:
    run = p.runs[0]
run.font.italic = True
run.font.size = Pt(28)
run.font.color.rgb = RGBColor(0x00, 0x40, 0x40)
p.alignment = PP_ALIGN.CENTER

out_path = os.path.join(os.getcwd(), 'Presentasi_Aplikasi_Bersihin.pptx')
prs.save(out_path)
print(f"Presentation saved successfully to: {out_path}")
